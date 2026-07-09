#!/usr/bin/env python3
"""
Convert a diagrams.net / draw.io file to an editable PowerPoint deck.

The converter focuses on visual equivalence for common flowchart elements:
rectangles, rounded rectangles, text, ellipses, diamonds, hexagons, cylinders,
notes, simple images, and orthogonal/straight connectors.

Usage:
    python drawio_to_pptx.py sample.drawio
    python drawio_to_pptx.py sample.drawio -o output.pptx --no-fix-mojibake
"""

from __future__ import annotations

import argparse
import base64
import html
import os
import re
import tempfile
import urllib.parse
import zlib
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Emu, Pt


EMU_PER_PX = 9525  # 1 px at 96 dpi = 0.75 pt = 9525 EMU


@dataclass(frozen=True)
class Point:
    x: float
    y: float


@dataclass(frozen=True)
class Box:
    x: float
    y: float
    w: float
    h: float


def px(value: float | int | str | None) -> Emu:
    return Emu(round(float(value or 0) * EMU_PER_PX))


def pt_from_px(value: float | int | str | None) -> Pt:
    return Pt(float(value or 0) * 0.75)


def parse_style(style: str | None) -> dict[str, str]:
    parsed: dict[str, str] = {}
    for part in (style or "").split(";"):
        if not part:
            continue
        if "=" in part:
            key, value = part.split("=", 1)
            parsed[key] = value
        else:
            parsed[part] = "1"
    return parsed


def rgb(value: str | None) -> RGBColor | None:
    if not value or value == "none":
        return None
    if value.startswith("#") and len(value) == 7:
        try:
            return RGBColor(int(value[1:3], 16), int(value[3:5], 16), int(value[5:7], 16))
        except ValueError:
            return None
    return None


def clean_text(value: str | None, fix_mojibake: bool = True) -> str:
    if not value:
        return ""
    text = html.unescape(value)
    text = re.sub(r"(?i)<br\s*/?>", "\n", text)
    text = re.sub(r"<[^>]+>", "", text)
    text = text.replace("\xa0", " ").strip()
    if fix_mojibake:
        text = repair_common_chinese_mojibake(text)
    return text


def repair_common_chinese_mojibake(text: str) -> str:
    """Repair text that was UTF-8 but got decoded as GBK/CP936."""
    if not text:
        return text
    try:
        repaired = text.encode("gbk").decode("utf-8")
    except UnicodeError:
        return text
    # Avoid changing already-correct Chinese. Correct strings usually fail the
    # transcode above; this score is a second guard for mixed ASCII/CJK labels.
    mojibake_marks = sum(text.count(ch) for ch in "鏅兘嵁紪绯荤粺娣卞寲瑙勫垯")
    if mojibake_marks or len(repaired) < len(text):
        return repaired
    return text


def decode_diagram_payload(payload: str) -> ET.Element:
    """Decode compressed draw.io diagram text into an mxGraphModel element."""
    raw = base64.b64decode(payload)
    try:
        inflated = zlib.decompress(raw, -15)
    except zlib.error:
        inflated = zlib.decompress(raw)
    xml_text = urllib.parse.unquote(inflated.decode("utf-8"))
    return ET.fromstring(xml_text)


def load_pages(drawio_path: Path) -> list[tuple[str, ET.Element]]:
    root = ET.parse(drawio_path).getroot()
    pages: list[tuple[str, ET.Element]] = []
    diagrams = root.findall("diagram") if root.tag == "mxfile" else []

    if root.tag == "mxGraphModel":
        pages.append(("Page-1", root))
    elif diagrams:
        for idx, diagram in enumerate(diagrams, start=1):
            name = diagram.get("name") or f"Page-{idx}"
            model = diagram.find("mxGraphModel")
            if model is None:
                payload = (diagram.text or "").strip()
                if not payload:
                    continue
                model = decode_diagram_payload(payload)
            pages.append((name, model))
    else:
        model = root.find(".//mxGraphModel")
        if model is not None:
            pages.append(("Page-1", model))

    if not pages:
        raise ValueError(f"No mxGraphModel page found in {drawio_path}")
    return pages


def get_geometry(cell: ET.Element) -> Box | None:
    geom = cell.find("mxGeometry")
    if geom is None:
        return None
    return Box(
        float(geom.get("x") or 0),
        float(geom.get("y") or 0),
        float(geom.get("width") or 0),
        float(geom.get("height") or 0),
    )


def build_cell_maps(model: ET.Element) -> tuple[dict[str, ET.Element], dict[str, Box]]:
    cells = {cell.get("id"): cell for cell in model.findall(".//mxCell") if cell.get("id")}
    boxes = {cell_id: box for cell_id, cell in cells.items() if (box := get_geometry(cell)) is not None}
    return cells, boxes


def absolute_box(cell_id: str, cells: dict[str, ET.Element], boxes: dict[str, Box]) -> Box:
    box = boxes.get(cell_id, Box(0, 0, 0, 0))
    parent_id = cells.get(cell_id).get("parent") if cell_id in cells else None
    if parent_id and parent_id in boxes and cells.get(parent_id, {}).get("vertex") == "1":
        parent = absolute_box(parent_id, cells, boxes)
        return Box(parent.x + box.x, parent.y + box.y, box.w, box.h)
    return box


def apply_fill_and_line(shape, style: dict[str, str]) -> None:
    fill_color = rgb(style.get("fillColor"))
    if fill_color is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        if "opacity" in style:
            shape.fill.transparency = max(0, min(100, 100 - float(style["opacity"]))) / 100

    stroke_color = rgb(style.get("strokeColor"))
    if stroke_color is None or style.get("strokeColor") == "none" or style.get("strokeWidth") == "0":
        shape.line.fill.background()
        return

    shape.line.color.rgb = stroke_color
    shape.line.width = pt_from_px(style.get("strokeWidth", "1"))
    if style.get("dashed") == "1":
        shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH


def apply_text(shape, text: str, style: dict[str, str]) -> None:
    if not text:
        return

    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0

    valign = style.get("verticalAlign", "middle")
    frame.vertical_anchor = {
        "top": MSO_ANCHOR.TOP,
        "bottom": MSO_ANCHOR.BOTTOM,
        "middle": MSO_ANCHOR.MIDDLE,
    }.get(valign, MSO_ANCHOR.MIDDLE)

    lines = text.splitlines() or [text]
    for idx, line in enumerate(lines):
        para = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        para.text = line
        para.alignment = {
            "left": PP_ALIGN.LEFT,
            "right": PP_ALIGN.RIGHT,
            "center": PP_ALIGN.CENTER,
        }.get(style.get("align", "center"), PP_ALIGN.CENTER)

        font = para.font
        font.name = style.get("fontFamily", "Microsoft YaHei")
        font.size = pt_from_px(style.get("fontSize", "18"))
        color = rgb(style.get("fontColor"))
        if color is not None:
            font.color.rgb = color

        font_style = int(style.get("fontStyle", "0") or 0)
        font.bold = bool(font_style & 1)
        font.italic = bool(font_style & 2)
        font.underline = bool(font_style & 4)


def shape_type_for(style: dict[str, str]) -> MSO_SHAPE:
    if "ellipse" in style:
        return MSO_SHAPE.OVAL
    if "rhombus" in style:
        return MSO_SHAPE.DIAMOND

    shape = style.get("shape", "")
    if shape == "cylinder":
        return MSO_SHAPE.CAN
    if shape == "hexagon":
        return MSO_SHAPE.HEXAGON
    if shape == "note":
        return MSO_SHAPE.FOLDED_CORNER
    if shape == "process":
        return MSO_SHAPE.FLOWCHART_PROCESS
    if style.get("rounded") == "1":
        return MSO_SHAPE.ROUNDED_RECTANGLE
    return MSO_SHAPE.RECTANGLE


def add_vertex(slide, cell: ET.Element, box: Box, fix_mojibake: bool) -> None:
    style = parse_style(cell.get("style"))
    text = clean_text(cell.get("value"), fix_mojibake)

    if style.get("shape") == "image" and style.get("image", "").startswith("data:image/"):
        add_data_uri_image(slide, style["image"], box)
        return

    if "text" in style:
        shape = slide.shapes.add_textbox(px(box.x), px(box.y), px(box.w), px(box.h))
        apply_text(shape, text, style)
    else:
        shape = slide.shapes.add_shape(shape_type_for(style), px(box.x), px(box.y), px(box.w), px(box.h))
        apply_fill_and_line(shape, style)
        apply_text(shape, text, style)

    if "rotation" in style:
        shape.rotation = float(style["rotation"])


def add_data_uri_image(slide, data_uri: str, box: Box) -> None:
    parts = data_uri.split(",", 1)
    if len(parts) < 2 or not parts[1].strip():
        return  # skip invalid/incomplete data URI
    header, encoded = parts
    suffix = ".png"
    if "jpeg" in header or "jpg" in header:
        suffix = ".jpg"
    try:
        img_bytes = base64.b64decode(encoded)
    except Exception:
        return  # skip undecodable image data
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp.write(img_bytes)
        temp_path = tmp.name
    try:
        slide.shapes.add_picture(temp_path, px(box.x), px(box.y), width=px(box.w), height=px(box.h))
    finally:
        try:
            os.unlink(temp_path)
        except OSError:
            pass


def edge_points(cell: ET.Element, cells: dict[str, ET.Element], boxes: dict[str, Box]) -> list[Point]:
    geom = cell.find("mxGeometry")
    points: list[Point] = []

    has_source_point = False
    has_target_point = False

    if geom is not None:
        source = geom.find("mxPoint[@as='sourcePoint']")
        if source is not None:
            points.append(Point(float(source.get("x") or 0), float(source.get("y") or 0)))
            has_source_point = True

        array = geom.find("Array[@as='points']")
        if array is not None:
            for item in array.findall("mxPoint"):
                points.append(Point(float(item.get("x") or 0), float(item.get("y") or 0)))

        target = geom.find("mxPoint[@as='targetPoint']")
        if target is not None:
            points.append(Point(float(target.get("x") or 0), float(target.get("y") or 0)))
            has_target_point = True

    # If no explicit source/target points, compute edge attachment from node geometry
    source_ref = cell.get("source")
    target_ref = cell.get("target")

    if not has_source_point and source_ref and source_ref in boxes:
        box = absolute_box(source_ref, cells, boxes)
        cx, cy = box.x + box.w / 2, box.y + box.h / 2
        if points:
            # Attach to the box edge closest to the first waypoint
            fx, fy = points[0].x, points[0].y
            points.insert(0, Point(_edge_attachment_x(box, cx, fx), _edge_attachment_y(box, cy, fy)))
        else:
            points.append(Point(cx, cy))

    if not has_target_point and target_ref and target_ref in boxes:
        box = absolute_box(target_ref, cells, boxes)
        cx, cy = box.x + box.w / 2, box.y + box.h / 2
        if points:
            # Attach to the box edge closest to the last waypoint
            lx, ly = points[-1].x, points[-1].y
            points.append(Point(_edge_attachment_x(box, cx, lx), _edge_attachment_y(box, cy, ly)))
        else:
            points.append(Point(cx, cy))

    return points


def _edge_attachment_x(box: Box, cx: float, target_x: float) -> float:
    """Return x coord on box edge closest to target_x."""
    if target_x < box.x:
        return box.x
    if target_x > box.x + box.w:
        return box.x + box.w
    return cx  # vertically aligned, use center


def _edge_attachment_y(box: Box, cy: float, target_y: float) -> float:
    """Return y coord on box edge closest to target_y."""
    if target_y < box.y:
        return box.y
    if target_y > box.y + box.h:
        return box.y + box.h
    return cy  # horizontally aligned, use center


def set_line_arrow(connector, at_end: bool = True) -> None:
    ln = connector.line._get_or_add_ln()
    tag = "a:tailEnd" if not at_end else "a:headEnd"
    for existing in list(ln):
        if existing.tag.endswith(tag.split(":", 1)[1]):
            ln.remove(existing)
    arrow = OxmlElement(tag)
    arrow.set("type", "triangle")
    arrow.set("w", "med")
    arrow.set("len", "med")
    ln.append(arrow)


def add_edge(slide, cell: ET.Element, cells: dict[str, ET.Element], boxes: dict[str, Box]) -> None:
    style = parse_style(cell.get("style"))
    points = edge_points(cell, cells, boxes)
    if len(points) < 2:
        return

    color = rgb(style.get("strokeColor")) or RGBColor(0, 0, 0)
    width = pt_from_px(style.get("strokeWidth", "1"))
    dashed = style.get("dashed") == "1"

    for idx, (start, end) in enumerate(zip(points, points[1:])):
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            px(start.x),
            px(start.y),
            px(end.x),
            px(end.y),
        )
        connector.line.color.rgb = color
        connector.line.width = width
        if dashed:
            connector.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        if idx == len(points) - 2 and style.get("endArrow", "none") != "none":
            set_line_arrow(connector, at_end=True)
        if idx == 0 and style.get("startArrow", "none") != "none":
            set_line_arrow(connector, at_end=False)


def page_size(model: ET.Element) -> tuple[float, float]:
    width = float(model.get("pageWidth") or 1280)
    height = float(model.get("pageHeight") or 720)
    return width, height


def iter_cells_in_draw_order(model: ET.Element) -> Iterable[ET.Element]:
    root = model.find("root")
    return root.findall("mxCell") if root is not None else model.findall(".//mxCell")


def convert(drawio_path: Path, pptx_path: Path, fix_mojibake: bool = True) -> None:
    pages = load_pages(drawio_path)
    pres = Presentation()
    pres.slide_width = px(page_size(pages[0][1])[0])
    pres.slide_height = px(page_size(pages[0][1])[1])
    blank = pres.slide_layouts[6]

    # Remove the starter slide that some template versions provide.
    while pres.slides:
        rel_id = pres.slides._sldIdLst[0].rId
        pres.part.drop_rel(rel_id)
        del pres.slides._sldIdLst[0]

    for _, model in pages:
        width, height = page_size(model)
        pres.slide_width = px(width)
        pres.slide_height = px(height)
        slide = pres.slides.add_slide(blank)
        cells, boxes = build_cell_maps(model)

        for cell in iter_cells_in_draw_order(model):
            cell_id = cell.get("id")
            if not cell_id or cell.get("visible") == "0":
                continue
            if cell.get("vertex") == "1":
                add_vertex(slide, cell, absolute_box(cell_id, cells, boxes), fix_mojibake)
            elif cell.get("edge") == "1":
                add_edge(slide, cell, cells, boxes)

    pptx_path.parent.mkdir(parents=True, exist_ok=True)
    pres.save(pptx_path)


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Convert draw.io .drawio files to editable .pptx files.")
    parser.add_argument("input", help="Input .drawio file")
    parser.add_argument(
        "-o",
        "--output",
        help="Output .pptx file; defaults to input basename + _converted.pptx",
    )
    parser.add_argument(
        "--no-fix-mojibake",
        action="store_true",
        help="Preserve text exactly as stored instead of repairing common UTF-8-as-GBK mojibake.",
    )
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    drawio_path = Path(args.input).expanduser().resolve()
    if not drawio_path.is_file():
        raise SystemExit(f"Input file not found: {drawio_path}")
    output = (
        Path(args.output).expanduser().resolve()
        if args.output
        else drawio_path.with_name(f"{drawio_path.stem}_converted.pptx")
    )
    convert(drawio_path, output, fix_mojibake=not args.no_fix_mojibake)
    print(f"Wrote {output}")


if __name__ == "__main__":
    main()
